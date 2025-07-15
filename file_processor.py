import os
import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import docx
import pandas as pd
import json
import yaml
import xml.etree.ElementTree as ET
from pptx import Presentation
import logging
import mimetypes
from pathlib import Path

logger = logging.getLogger(__name__)

class FileProcessor:
    def __init__(self):
        self.processors = {
            'pdf': self._process_pdf,
            'doc': self._process_doc,
            'docx': self._process_docx,
            'txt': self._process_txt,
            'rtf': self._process_txt,
            'md': self._process_md,
            'html': self._process_html,
            'htm': self._process_html,
            'json': self._process_json,
            'xml': self._process_xml,
            'yaml': self._process_yaml,
            'yml': self._process_yaml,
            'csv': self._process_csv,
            'xls': self._process_excel,
            'xlsx': self._process_excel,
            'jpg': self._process_image,
            'jpeg': self._process_image,
            'png': self._process_image,
            'gif': self._process_image,
            'bmp': self._process_image,
            'tiff': self._process_image,
            'svg': self._process_vector,
            'mp3': self._process_audio,
            'wav': self._process_audio,
            'flac': self._process_audio,
            'ogg': self._process_audio,
            'm4a': self._process_audio,
            'mp4': self._process_video,
            'avi': self._process_video,
            'mov': self._process_video,
            'mkv': self._process_video,
            'zip': self._process_archive,
            'rar': self._process_archive,
            '7z': self._process_archive,
            'tar': self._process_archive,
            'gz': self._process_archive,
            'exe': self._process_executable,
            'dll': self._process_executable,
            'sh': self._process_script,
            'py': self._process_script,
            'js': self._process_script,
            'css': self._process_html,
            'log': self._process_log,
            'epub': self._process_epub,
            'iso': self._process_disk_image,
            'dmg': self._process_disk_image,
            'ppt': self._process_ppt,
            'pptx': self._process_ppt,
            'ppk': self._process_key_file,
            'pem': self._process_key_file,
            'key': self._process_key_file,
        }

    def _process_md(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Failed to process markdown file: {e}")
            return None

    def extract_content(self, filepath, filename):
        try:
            file_ext = self.get_file_extension(filename)
            if file_ext in self.processors:
                return self.processors[file_ext](filepath)
            else:
                logger.warning(f"Unsupported file type: {file_ext}")
                return None
        except Exception as e:
            logger.error(f"Error extracting content from {filename}: {str(e)}")
            return None

    def get_file_extension(self, filename):
        return filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

    def get_file_type(self, filename):
        ext = self.get_file_extension(filename)
        if ext in ['pdf', 'doc', 'docx', 'txt', 'rtf', 'md']:
            return 'document'
        elif ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff']:
            return 'image'
        elif ext in ['xlsx', 'xls', 'csv']:
            return 'spreadsheet'
        elif ext in ['ppt', 'pptx']:
            return 'presentation'
        elif ext in ['ppk', 'pem', 'key']:
            return 'key_file'
        elif ext in ['json', 'xml', 'yaml', 'yml']:
            return 'data_file'
        elif ext in ['html', 'htm']:
            return 'web_page'
        else:
            return 'unknown'

    def _process_pdf(self, filepath):
        try:
            doc = fitz.open(filepath)
            text = ""
            for page in doc:
                text += page.get_text()
            return text.strip()
        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            return None

    def _process_docx(self, filepath):
        try:
            doc = docx.Document(filepath)
            return '\n'.join(p.text for p in doc.paragraphs)
        except Exception as e:
            logger.error(f"Error processing DOCX: {str(e)}")
            return None

    def _process_doc(self, filepath):
        try:
            import docx2txt
            return docx2txt.process(filepath)
        except ImportError:
            logger.warning("docx2txt not installed, cannot process .doc files")
            return "DOC file detected but docx2txt not installed"
        except Exception as e:
            logger.error(f"Error processing DOC: {str(e)}")
            return None

    def _process_txt(self, filepath):
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            try:
                with open(filepath, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Error processing TXT with latin-1: {str(e)}")
                return None
        except Exception as e:
            logger.error(f"Error processing TXT: {str(e)}")
            return None

    def _process_image(self, filepath):
        try:
            img = Image.open(filepath)
            return pytesseract.image_to_string(img).strip()
        except Exception as e:
            logger.error(f"Error processing image: {str(e)}")
            return None

    def _process_excel(self, filepath):
        try:
            df = pd.read_excel(filepath, sheet_name=None)
            return '\n\n'.join(f"Sheet: {name}\n{data.to_string()}" for name, data in df.items())
        except Exception as e:
            logger.error(f"Error processing Excel: {str(e)}")
            return None

    def _process_csv(self, filepath):
        try:
            df = pd.read_csv(filepath)
            return df.to_string()
        except Exception as e:
            logger.error(f"Error processing CSV: {str(e)}")
            return None

    def _process_json(self, filepath):
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return json.dumps(json.load(f), indent=2)
        except Exception as e:
            logger.error(f"Error processing JSON: {str(e)}")
            return None

    def _process_xml(self, filepath):
        try:
            tree = ET.parse(filepath)
            return ET.tostring(tree.getroot(), encoding='unicode')
        except Exception as e:
            logger.error(f"Error processing XML: {str(e)}")
            return None

    def _process_yaml(self, filepath):
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return yaml.dump(yaml.safe_load(f))
        except Exception as e:
            logger.error(f"Error processing YAML: {str(e)}")
            return None

    def _process_html(self, filepath):
        try:
            from bs4 import BeautifulSoup
            with open(filepath, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, 'html.parser')
            return soup.get_text()
        except ImportError:
            logger.warning("BeautifulSoup not installed, processing HTML as plain text")
            return self._process_txt(filepath)
        except Exception as e:
            logger.error(f"Error processing HTML: {str(e)}")
            return None

    # Placeholders for the rest:
    def _process_vector(self, filepath): return "Vector file processing not implemented"
    def _process_audio(self, filepath): return "Audio file processing not implemented"
    def _process_video(self, filepath): return "Video file processing not implemented"
    def _process_archive(self, filepath): return "Archive file processing not implemented"
    def _process_executable(self, filepath): return "Executable file detected"
    def _process_script(self, filepath): return self._process_txt(filepath)
    def _process_log(self, filepath): return self._process_txt(filepath)
    def _process_epub(self, filepath): return "EPUB file processing not implemented"
    def _process_disk_image(self, filepath): return "Disk image file processing not implemented"
    def _process_ppt(self, filepath):
        try:
            prs = Presentation(filepath)
            return '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
        except Exception as e:
            logger.error(f"Error processing PPT/PPTX: {str(e)}")
            return None
    def _process_key_file(self, filepath): return "Key file detected; processing skipped for security"


